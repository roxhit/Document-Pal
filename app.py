import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings import OpenAIEmbeddings, HuggingFaceInstructEmbeddings
from langchain.vectorstores import FAISS
from langchain.chat_models import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from templates import css, bot_template, user_template
from langchain.llms import HuggingFaceHub
from io import BytesIO
import docx
from pptx import Presentation

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            text += page.extract_text()
    return text

def extract_text_from_docx(docx_file):
    docx_content = docx_file.read()  # Read the bytes content of the file
    doc = docx.Document(BytesIO(docx_content))  # Convert to file-like object
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    return chunks

def get_text_from_files(files):
    text = ""
    for file in files:
        if file.name.endswith('.pdf'):
            text += get_pdf_text(file)
        elif file.name.endswith('.docx'):
            text += extract_text_from_docx(file)
        elif file.name.endswith('.pptx'):
            text += extract_text_from_pptx(file)
    return text

def get_vectorstore(text_chunks):
    embeddings = OpenAIEmbeddings()
    # embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    vectorstore = FAISS.from_texts(texts=text_chunks, embedding=embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    llm = ChatOpenAI()
    # llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})

    memory = ConversationBufferMemory(
        memory_key='chat_history', return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm=llm,
        retriever=vectorstore.as_retriever(),
        memory=memory
    )
    return conversation_chain

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']

    for i, message in enumerate(st.session_state.chat_history):
        if i % 2 == 0:
            st.write(user_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)
        else:
            st.write(bot_template.replace(
                "{{MSG}}", message.content), unsafe_allow_html=True)

load_dotenv()
st.set_page_config(page_title="Chat with multiple PDFs", page_icon=":books:")
st.write(css, unsafe_allow_html=True)

if "conversation" not in st.session_state:
    st.session_state.conversation = None
if "chat_history" not in st.session_state:
    st.session_state.chat_history = None

st.header("Document-Pal :books:")
user_question = st.text_input("Ask a question about your documents:")
if user_question:
    handle_userinput(user_question)

with st.sidebar:
    st.subheader("Your documents")
    selected_file_type = st.radio("Select file type:", ["PDF", "DOCX", "PPTX"])

    if selected_file_type == "PDF":
        uploaded_files = st.file_uploader(
            "Upload your PDFs here", accept_multiple_files=True, key="pdf_uploader")
    elif selected_file_type == "DOCX":
        uploaded_files = st.file_uploader(
            "Upload your DOCX files here", accept_multiple_files=True, key="docx_uploader")
    elif selected_file_type == "PPTX":
        uploaded_files = st.file_uploader(
            "Upload your PPTX files here", accept_multiple_files=True, key="pptx_uploader")

if st.button("Process"):
    with st.spinner("Processing"):
        if uploaded_files:
            raw_text = get_text_from_files(uploaded_files)
            text_chunks = get_text_chunks(raw_text)
            vectorstore = get_vectorstore(text_chunks)
            conversation_chain = get_conversation_chain(vectorstore)
            st.session_state.conversation = conversation_chain
